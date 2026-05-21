from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import copy

# ── 색상 팔레트 ──────────────────────────────────────────────
C_DARK   = RGBColor(0x1E, 0x29, 0x3B)   # 슬레이트 900
C_BLUE   = RGBColor(0x1D, 0x4E, 0xD8)   # 블루 700
C_TEAL   = RGBColor(0x0F, 0x76, 0x6E)   # 틸 700
C_RED    = RGBColor(0xBE, 0x12, 0x3C)   # 로즈 700
C_GRAY   = RGBColor(0x64, 0x74, 0x8B)   # 슬레이트 500
C_LGRAY  = RGBColor(0xF1, 0xF5, 0xF9)   # 슬레이트 100
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_AMBER  = RGBColor(0xD9, 0x77, 0x06)   # 앰버 600

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # 완전 빈 레이아웃


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, x, y, w, h, text, size=14, bold=False, color=C_DARK,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def header_bar(slide, title, subtitle=""):
    add_rect(slide, 0, 0, 13.33, 1.1, fill=C_DARK)
    add_textbox(slide, 0.4, 0.15, 10, 0.5, title,
                size=24, bold=True, color=C_WHITE)
    if subtitle:
        add_textbox(slide, 0.4, 0.65, 10, 0.35, subtitle,
                    size=13, color=RGBColor(0x94, 0xA3, 0xB8))


# ══════════════════════════════════════════════════════════════
# 슬라이드 1 — 표지
# ══════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
add_rect(s1, 0, 0, 13.33, 7.5, fill=C_DARK)
add_rect(s1, 0, 5.8, 13.33, 1.7, fill=C_BLUE)

add_textbox(s1, 1.0, 1.5, 11, 1.0,
            "SW FMEA 자동화 시스템",
            size=36, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
add_textbox(s1, 1.0, 2.6, 11, 0.7,
            "인턴 과제(SW Automation Hub) 대비 차별점",
            size=20, color=RGBColor(0x93, 0xC5, 0xFD), align=PP_ALIGN.LEFT)
add_textbox(s1, 1.0, 3.4, 11, 0.5,
            "전동화SW설계1팀  ·  우정호",
            size=14, color=C_GRAY, align=PP_ALIGN.LEFT)
add_textbox(s1, 1.0, 6.0, 4, 0.4,
            "2026. 05. 21",
            size=13, color=C_WHITE)

# ══════════════════════════════════════════════════════════════
# 슬라이드 2 — 핵심 한 줄 요약
# ══════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
add_rect(s2, 0, 0, 13.33, 7.5, fill=C_LGRAY)
header_bar(s2, "핵심 차이", "두 시스템이 해결하는 문제의 수준이 다르다")

# 두 박스
for bx, title, desc, color in [
    (0.4, "인턴 과제\n(SW Automation Hub)",
     "Interface xlsx / Drawio / 자연어\n→ FMEA 틀·SADS 문서·다이어그램 파일 생성\n\nS / O / D 점수는 여전히 수기 작성\n\n\"사람이 할 일을 줄여주는 보조 도구\"",
     C_RED),
    (6.9, "본 시스템\n(FMEA 자동화 플랫폼)",
     "ARXML 파일\n→ AI가 신호 의미를 해석 → S/O/D 자동 산출\n→ 가이드라인 기준으로 자동 보정\n\n\"사람 대신 판단하는 분석 플랫폼\"",
     C_TEAL),
]:
    add_rect(s2, bx, 1.35, 5.9, 5.5, fill=C_WHITE,
             line=color, line_w=Pt(2))
    add_rect(s2, bx, 1.35, 5.9, 0.65, fill=color)
    add_textbox(s2, bx+0.15, 1.38, 5.6, 0.6, title,
                size=14, bold=True, color=C_WHITE)
    add_textbox(s2, bx+0.2, 2.1, 5.5, 4.5, desc,
                size=13, color=C_DARK)

# 화살표 대신 VS 텍스트
add_textbox(s2, 6.15, 3.7, 1.0, 0.6, "VS",
            size=22, bold=True, color=C_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# 슬라이드 3 — 항목별 비교표
# ══════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
add_rect(s3, 0, 0, 13.33, 7.5, fill=C_LGRAY)
header_bar(s3, "항목별 비교")

rows = [
    ("항목",                     "인턴 과제",                       "본 시스템",                         True),
    ("S/O/D 점수 생성",          "수기 작성",                       "AI 자동 생성",                      False),
    ("입력 파일",                "Interface.xlsx 수동 준비",        "ARXML 직접 파싱",                   False),
    ("품질 기준 적용",           "없음",                            "SL FMEA Guideline v4.2 자동 보정",  False),
    ("플랫폼",                   "로컬 실행 파일 (.exe)",           "웹 앱 (브라우저·클라우드)",          False),
    ("데이터 저장",              "파일 시스템 (Excel/docx)",        "DB (Supabase)",                     False),
    ("프로젝트 관리",            "단일 차종",                       "다중 차종 동시 관리 (JG1/LQ2/NQ6e)", False),
    ("분석·통계",                "없음",                            "RPN 대시보드·프로젝트 비교",         False),
    ("개념 FMEA",                "없음",                            "Concept 단계 FMEA 자동 생성",        False),
]

col_w = [3.2, 4.2, 4.8]
col_x = [0.25, 3.55, 7.85]
row_h = 0.52
start_y = 1.2

for ri, (item, intern, ours, is_hdr) in enumerate(rows):
    y = start_y + ri * row_h
    bg = C_DARK if is_hdr else (C_WHITE if ri % 2 == 1 else RGBColor(0xF8, 0xFA, 0xFC))
    for ci, (text, cx, cw) in enumerate(zip([item, intern, ours], col_x, col_w)):
        add_rect(s3, cx, y, cw, row_h - 0.03,
                 fill=bg, line=RGBColor(0xE2, 0xE8, 0xF0), line_w=Pt(0.5))
        fc = C_WHITE if is_hdr else (C_TEAL if ci == 2 and not is_hdr and text not in ("없음", "수기 작성", "단일 차종", "파일 시스템 (Excel/docx)", "로컬 실행 파일 (.exe)", "Interface.xlsx 수동 준비") else (C_RED if text in ("없음", "수기 작성") else C_DARK))
        add_textbox(s3, cx+0.1, y+0.08, cw-0.15, row_h-0.1, text,
                    size=11 if not is_hdr else 12,
                    bold=is_hdr, color=fc)

# ══════════════════════════════════════════════════════════════
# 슬라이드 4 — 차별점 ① S/O/D 자동 생성
# ══════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
add_rect(s4, 0, 0, 13.33, 7.5, fill=C_LGRAY)
header_bar(s4, "차별점 ①  S/O/D 자동 생성", "가장 큰 공수 절감 포인트")

# 인턴 흐름
add_rect(s4, 0.4, 1.4, 2.5, 0.6, fill=C_RED)
add_textbox(s4, 0.4, 1.4, 2.5, 0.6, "Interface.xlsx", size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_textbox(s4, 3.1, 1.55, 0.5, 0.4, "→", size=16, bold=True, color=C_GRAY)
add_rect(s4, 3.7, 1.4, 2.5, 0.6, fill=RGBColor(0xFE, 0xF2, 0xF2), line=C_RED, line_w=Pt(1))
add_textbox(s4, 3.7, 1.4, 2.5, 0.6, "FMEA 틀 (행/열)", size=12, color=C_DARK, align=PP_ALIGN.CENTER)
add_textbox(s4, 6.3, 1.55, 0.5, 0.4, "→", size=16, bold=True, color=C_GRAY)
add_rect(s4, 6.9, 1.4, 2.8, 0.6, fill=RGBColor(0xFE, 0xF2, 0xF2), line=C_RED, line_w=Pt(1))
add_textbox(s4, 6.9, 1.4, 2.8, 0.6, "S/O/D  ← 수기 작성", size=12, color=C_RED, align=PP_ALIGN.CENTER)
add_textbox(s4, 0.4, 1.1, 4.0, 0.3, "인턴 과제", size=11, bold=True, color=C_RED)

# 본 시스템 흐름
add_rect(s4, 0.4, 3.3, 2.5, 0.6, fill=C_TEAL)
add_textbox(s4, 0.4, 3.3, 2.5, 0.6, "ARXML 파일", size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_textbox(s4, 3.1, 3.45, 0.5, 0.4, "→", size=16, bold=True, color=C_GRAY)
add_rect(s4, 3.7, 3.3, 2.5, 0.6, fill=RGBColor(0xF0, 0xFD, 0xFA), line=C_TEAL, line_w=Pt(1))
add_textbox(s4, 3.7, 3.3, 2.5, 0.6, "신호 자동 파싱", size=12, color=C_DARK, align=PP_ALIGN.CENTER)
add_textbox(s4, 6.3, 3.45, 0.5, 0.4, "→", size=16, bold=True, color=C_GRAY)
add_rect(s4, 6.9, 3.3, 2.5, 0.6, fill=RGBColor(0xF0, 0xFD, 0xFA), line=C_TEAL, line_w=Pt(1))
add_textbox(s4, 6.9, 3.3, 2.5, 0.6, "Claude AI 분석", size=12, color=C_DARK, align=PP_ALIGN.CENTER)
add_textbox(s4, 9.5, 3.45, 0.5, 0.4, "→", size=16, bold=True, color=C_GRAY)
add_rect(s4, 10.1, 3.3, 2.8, 0.6, fill=C_TEAL)
add_textbox(s4, 10.1, 3.3, 2.8, 0.6, "S/O/D 자동 산출", size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_textbox(s4, 0.4, 3.0, 4.0, 0.3, "본 시스템", size=11, bold=True, color=C_TEAL)

# 추가 설명
add_rect(s4, 0.4, 4.3, 12.5, 2.7, fill=C_WHITE, line=C_TEAL, line_w=Pt(1))
add_textbox(s4, 0.7, 4.45, 12.0, 2.4,
            "• Interface.xlsx 준비 단계 자체가 생략됨 (ARXML에서 직접 추출)\n"
            "• AI가 신호 이름·타입·범위를 해석해 Severity / Occurrence / Detection 점수 결정\n"
            "• SL SW FMEA Guideline v4.2 규칙을 내재화 → 가이드라인 기준 자동 보정\n"
            "• 담당자는 AI 결과 검토·수정만 하면 됨",
            size=13, color=C_DARK)

# ══════════════════════════════════════════════════════════════
# 슬라이드 5 — 차별점 ② 플랫폼·관리
# ══════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
add_rect(s5, 0, 0, 13.33, 7.5, fill=C_LGRAY)
header_bar(s5, "차별점 ②  플랫폼 및 데이터 관리", "로컬 파일 vs 클라우드 DB")

cards = [
    ("로컬 실행 파일\n(인턴 과제)", "• Windows PC에 설치 필요\n• 결과물이 로컬 Excel/docx\n• 버전 관리 없음\n• 1인 단독 사용", C_RED),
    ("웹 앱 + DB\n(본 시스템)", "• 브라우저로 즉시 접근\n• Supabase DB에 항목 저장\n• 다중 프로젝트 동시 관리\n• 팀원 공유 가능", C_TEAL),
    ("분석 기능\n(본 시스템 전용)", "• RPN 대시보드 (위험도 시각화)\n• 두 프로젝트 간 비교 분석\n• Concept 단계 FMEA 생성\n• 항목별 필터·검색", C_BLUE),
]
for ci, (title, body, color) in enumerate(cards):
    x = 0.4 + ci * 4.3
    add_rect(s5, x, 1.35, 4.0, 5.6, fill=C_WHITE, line=color, line_w=Pt(2))
    add_rect(s5, x, 1.35, 4.0, 0.65, fill=color)
    add_textbox(s5, x+0.15, 1.38, 3.7, 0.6, title,
                size=13, bold=True, color=C_WHITE)
    add_textbox(s5, x+0.2, 2.15, 3.6, 4.5, body,
                size=12, color=C_DARK)

# ══════════════════════════════════════════════════════════════
# 슬라이드 6 — 결론
# ══════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
add_rect(s6, 0, 0, 13.33, 7.5, fill=C_DARK)
add_rect(s6, 0, 5.5, 13.33, 2.0, fill=C_BLUE)

add_textbox(s6, 1.0, 1.0, 11, 0.8,
            "결론",
            size=28, bold=True, color=RGBColor(0x93, 0xC5, 0xFD))

add_textbox(s6, 1.0, 2.0, 11.3, 1.0,
            "인턴 과제는 \"사람이 할 일을 줄여주는 보조 도구\"",
            size=18, color=RGBColor(0xFC, 0xA5, 0xA5))

add_textbox(s6, 1.0, 3.1, 11.3, 1.0,
            "본 시스템은 \"사람 대신 판단하는 FMEA 분석 플랫폼\"",
            size=18, bold=True, color=C_WHITE)

details = [
    "① ARXML → S/O/D 자동 생성   (내용 생성, 틀 생성 X)",
    "② SL FMEA Guideline v4.2 기반 자동 보정   (품질 기준 내재화)",
    "③ 웹 기반 + Supabase DB   (협업·접근성)",
    "④ 다중 차종 관리 + RPN 대시보드 + 프로젝트 비교   (분석 플랫폼)",
]
add_textbox(s6, 1.2, 5.6, 11, 1.6,
            "\n".join(details),
            size=12, color=C_WHITE)

# ── 저장 ──────────────────────────────────────────────────────
out = "FMEA_자동화_차별점.pptx"
prs.save(out)
print(f"저장 완료: {out}")
